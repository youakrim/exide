#!/usr/bin/python
#-*- coding: utf-8 -*-
import getopt
import os

import sys

import pptx


# On cherche à récuperer les noms des fichiers d'entrée et de sortie
import pptx_parser


def main(argv):
    # We try to get the input and output files names
    inputfile = ''
    outputfile = ''
    try:
        opts, args = getopt.getopt(argv, "hi:o:", ["ifile=", "ofile="])
        '''print("Opts : ")
        print(opts)
        print("Args : ")
        print(args)'''
    except getopt.GetoptError:
        print('test.py -i <inputfile> -o <outputfile>')
        sys.exit(2)
    for opt, arg in opts:
        if opt == '-h':
            print('test.py -i <inputfile> -o <outputfile>')
            sys.exit()
        elif opt in ("-i", "--ifile"):
            inputfile = arg
        elif opt in ("-o", "--ofile"):
            outputfile = arg
    #print('Input file is "', inputfile)
    #print('Output file is "', outputfile)

    # On récupère les noms de fichiers et extensions
    # We extract the filename and extension of both files
    inputfile_name, inputfile_extension = os.path.splitext(inputfile)
    inputfile_directory = os.path.dirname(inputfile)
    #print("L'extension du fichier d'entrée est", inputfile_extension)
    # On vérifie que l'extension du fichier d'entrée (inputfile) est connue
    # We check if the inputfile extension is recognized
    extensions_connues = [".pptx"]

    if inputfile_extension not in extensions_connues:
        sys.exit("Erreur : Type de fichier non supporté \n \t Les extensions supportées sont ", extensions_connues)

    # Si odp alors on entre dans l'archive
    # If the input file extension is ".odp", we unzip the archive
    if inputfile_extension == ".pptx":
        prs = pptx.Presentation(inputfile)
        for slide in prs.slides:
            new_title = "[title]"+slide.shapes.title.text + "[/title]"
            slide.shapes.title.text = new_title
        prs.save(outputfile)

if __name__ == "__main__":
    main(sys.argv[1:])
